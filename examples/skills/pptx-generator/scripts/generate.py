#!/usr/bin/env python3
"""Generate /workspace/output/deck.pptx from an outline.md staged by the thread.

The platform's `skill_run` tool stages every chat-uploaded attachment into
/workspace/output/<filename> before this script runs. We look for the
first markdown file matching outline.md / *.outline.md and parse it into
slides; if nothing matches, we fall back to a small built-in demo deck so
a cold invocation still produces a visible artifact.

Output: /workspace/output/deck.pptx, harvested by the sandbox and
returned to the thread.
"""

import glob
import sys
from dataclasses import dataclass, field
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError as e:
    sys.stderr.write(f"python-pptx import failed: {e}\n")
    sys.exit(1)

WORKDIR = Path("/workspace/output")
OUT_PATH = WORKDIR / "deck.pptx"


@dataclass
class Slide:
    title: str
    bullets: list[str] = field(default_factory=list)
    notes: list[str] = field(default_factory=list)


@dataclass
class Outline:
    title: str
    slides: list[Slide] = field(default_factory=list)


def find_outline_file() -> Path | None:
    candidates: list[Path] = []
    for pattern in ("outline.md", "*.outline.md", "outline.markdown"):
        candidates.extend(Path(p) for p in glob.glob(str(WORKDIR / pattern)))
    candidates = sorted(set(candidates))
    return candidates[0] if candidates else None


def parse_outline(text: str) -> Outline:
    """Tiny markdown subset parser — H1=title, H2=slide, list items=bullets,
    other prose between bullets and next H2=speaker notes for the slide."""
    outline = Outline(title="Untitled deck")
    title_seen = False
    current: Slide | None = None

    for raw_line in text.splitlines():
        line = raw_line.rstrip()
        stripped = line.strip()
        if not stripped:
            continue

        if stripped.startswith("# ") and not stripped.startswith("## "):
            if not title_seen:
                outline.title = stripped[2:].strip() or outline.title
                title_seen = True
            continue

        if stripped.startswith("## "):
            current = Slide(
                title=stripped[3:].strip() or f"Slide {len(outline.slides) + 1}"
            )
            outline.slides.append(current)
            continue

        if current is None:
            # Prose before any ## heading — promote into a leading slide so
            # we don't drop content.
            current = Slide(title=outline.title)
            outline.slides.append(current)

        if stripped.startswith(("-", "*", "+")):
            current.bullets.append(stripped.lstrip("-*+ ").strip())
        else:
            current.notes.append(stripped)

    if not outline.slides:
        outline.slides.append(
            Slide(title=outline.title, notes=["(outline had no slide-level headings)"])
        )
    return outline


def demo_outline() -> Outline:
    return Outline(
        title="Sample deck",
        slides=[
            Slide(
                title="What this skill does",
                bullets=[
                    "Turns a markdown outline into a .pptx",
                    "Title slide + one slide per `## Heading`",
                    "Bullets from list items, prose becomes speaker notes",
                ],
                notes=["Demo fallback — attach outline.md to drive real content."],
            ),
            Slide(
                title="How to use it",
                bullets=[
                    "Attach outline.md to the thread",
                    "Ask the agent for a deck",
                    "Receive deck.pptx as an attachment",
                ],
            ),
            Slide(
                title="Extending",
                bullets=[
                    "Edit scripts/generate.py to add images, charts, themes",
                    "Add new fonts or layouts via python-pptx",
                ],
                notes=[
                    "The script intentionally stays simple so the bundle is readable."
                ],
            ),
        ],
    )


def build_deck(outline: Outline, out_path: Path) -> None:
    prs = Presentation()

    # Title slide — layout 0 is "Title Slide" in the default template.
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = outline.title
    if title_slide.placeholders and len(title_slide.placeholders) > 1:
        subtitle = title_slide.placeholders[1]
        subtitle.text = f"{len(outline.slides)} slide(s)"

    # Layout 1 is "Title and Content" — bullets in a single placeholder.
    for slide_spec in outline.slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = slide_spec.title

        body_placeholder = None
        for shape in slide.placeholders:
            # idx 1 is the content placeholder on the default Title+Content
            # layout. Title is idx 0.
            if shape.placeholder_format.idx == 1:
                body_placeholder = shape
                break

        if body_placeholder is not None and slide_spec.bullets:
            tf = body_placeholder.text_frame
            tf.text = slide_spec.bullets[0]
            for bullet in slide_spec.bullets[1:]:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0
            # Bump default font a touch for readability when the deck is
            # projected — keeps the layout but lifts ~18pt → ~24pt.
            for paragraph in tf.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(24)

        if slide_spec.notes:
            notes_tf = slide.notes_slide.notes_text_frame
            notes_tf.text = " ".join(slide_spec.notes)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))


def main() -> int:
    outline_path = find_outline_file()
    if outline_path is None:
        outline = demo_outline()
        source = "(built-in demo)"
    else:
        try:
            text = outline_path.read_text(encoding="utf-8")
        except OSError as exc:
            sys.stderr.write(f"could not read {outline_path}: {exc}\n")
            return 1
        outline = parse_outline(text)
        source = outline_path.name

    build_deck(outline, OUT_PATH)

    sys.stdout.write(
        f"Wrote {OUT_PATH} — {len(outline.slides)} content slide(s) from {source}.\n"
    )
    if outline.slides:
        sys.stdout.write("Slide titles:\n")
        for i, slide in enumerate(outline.slides, start=1):
            sys.stdout.write(f"  {i}. {slide.title}\n")
    return 0


if __name__ == "__main__":
    sys.exit(main())
