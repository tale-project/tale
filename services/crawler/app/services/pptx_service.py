"""
PPTX Service for PowerPoint generation and analysis.

Handles:
- PPTX template analysis (extract structure, placeholders, charts)
- PPTX generation from template + content
"""

import logging
from io import BytesIO
from typing import Optional, List, Dict, Any

import httpx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.chart import XL_CHART_TYPE

logger = logging.getLogger(__name__)


class PptxService:
    """Service for analyzing and generating PPTX documents."""

    def __init__(self):
        self._http_client: Optional[httpx.AsyncClient] = None

    async def _get_http_client(self) -> httpx.AsyncClient:
        """Get or create HTTP client."""
        if self._http_client is None:
            self._http_client = httpx.AsyncClient(timeout=60.0)
        return self._http_client

    async def cleanup(self):
        """Cleanup resources."""
        if self._http_client:
            await self._http_client.aclose()
            self._http_client = None

    async def download_file(self, url: str) -> bytes:
        """Download a file from URL."""
        client = await self._get_http_client()
        response = await client.get(url)
        response.raise_for_status()
        return response.content

    def _get_chart_type_name(self, chart_type: int) -> str:
        """Convert chart type enum to readable name."""
        chart_type_names = {
            XL_CHART_TYPE.AREA: "area",
            XL_CHART_TYPE.AREA_STACKED: "area_stacked",
            XL_CHART_TYPE.BAR_CLUSTERED: "bar_clustered",
            XL_CHART_TYPE.BAR_STACKED: "bar_stacked",
            XL_CHART_TYPE.COLUMN_CLUSTERED: "column_clustered",
            XL_CHART_TYPE.COLUMN_STACKED: "column_stacked",
            XL_CHART_TYPE.LINE: "line",
            XL_CHART_TYPE.LINE_MARKERS: "line_markers",
            XL_CHART_TYPE.PIE: "pie",
            XL_CHART_TYPE.DOUGHNUT: "doughnut",
            XL_CHART_TYPE.RADAR: "radar",
            XL_CHART_TYPE.XY_SCATTER: "scatter",
        }
        return chart_type_names.get(chart_type, f"unknown_{chart_type}")

    def _get_shape_text(self, shape) -> str:
        """Safely get text from a shape."""
        try:
            if shape.has_text_frame:
                return shape.text_frame.text
        except Exception:
            pass
        return ""

    def _parse_hex_color(self, hex_color: Optional[str]):
        """Parse a hex color string to RGBColor."""
        from pptx.dml.color import RGBColor

        if not hex_color:
            return None
        try:
            hex_color = hex_color.lstrip("#")
            if len(hex_color) == 6:
                return RGBColor(
                    int(hex_color[0:2], 16),
                    int(hex_color[2:4], 16),
                    int(hex_color[4:6], 16)
                )
        except Exception:
            pass
        return None

    async def analyze_pptx_template(
        self,
        template_url: Optional[str] = None,
        template_bytes: Optional[bytes] = None,
        template_base64: Optional[str] = None,
    ) -> Dict[str, Any]:
        """
        Analyze a PPTX template and extract its FULL structure and content.

        Returns:
            Dict with slide count, layouts, and full content of each slide.
        """
        import base64

        if template_url:
            template_bytes = await self.download_file(template_url)
        elif template_base64:
            template_bytes = base64.b64decode(template_base64)

        if not template_bytes:
            raise ValueError("Either template_url, template_bytes, or template_base64 must be provided")

        prs = Presentation(BytesIO(template_bytes))
        slides_info = []

        for idx, slide in enumerate(prs.slides):
            slide_info = self._extract_slide_info(idx, slide)
            slides_info.append(slide_info)

        layouts = []
        try:
            for layout in prs.slide_layouts:
                layouts.append(layout.name)
        except Exception:
            pass

        branding = self._extract_branding(prs)

        return {
            "success": True,
            "slideCount": len(prs.slides),
            "slides": slides_info,
            "availableLayouts": layouts,
            "branding": branding,
        }

    def _extract_slide_info(self, idx: int, slide) -> Dict[str, Any]:
        """Extract information from a single slide."""
        slide_info = {
            "slideNumber": idx + 1,
            "layoutName": slide.slide_layout.name if slide.slide_layout else "Unknown",
            "title": None,
            "subtitle": None,
            "textContent": [],
            "tables": [],
            "charts": [],
            "images": [],
        }

        for shape in slide.shapes:
            self._extract_shape_info(shape, slide_info)

        return slide_info

    def _extract_shape_info(self, shape, slide_info: Dict[str, Any]) -> None:
        """Extract information from a shape and add to slide_info."""
        # Extract title/subtitle from placeholders
        if shape.is_placeholder:
            try:
                ph_type = shape.placeholder_format.type.name if shape.placeholder_format.type else ""
                text = self._get_shape_text(shape)
                if ph_type == "TITLE" or ph_type == "CENTER_TITLE":
                    slide_info["title"] = text
                elif ph_type == "SUBTITLE":
                    slide_info["subtitle"] = text
            except Exception:
                pass

        # Extract all text content
        if shape.has_text_frame:
            text = self._get_shape_text(shape)
            if text and text.strip():
                slide_info["textContent"].append({
                    "text": text,
                    "isPlaceholder": shape.is_placeholder,
                })

        # Extract table data
        if shape.has_table:
            try:
                table = shape.table
                table_data = {
                    "rowCount": len(table.rows),
                    "columnCount": len(table.columns),
                    "headers": [],
                    "rows": [],
                }
                for row_idx, row in enumerate(table.rows):
                    row_data = []
                    for cell in row.cells:
                        row_data.append(cell.text if cell.text else "")
                    if row_idx == 0:
                        table_data["headers"] = row_data
                    else:
                        table_data["rows"].append(row_data)
                slide_info["tables"].append(table_data)
            except Exception as e:
                logger.warning(f"Error extracting table data: {e}")

        # Extract chart info
        if shape.has_chart:
            try:
                chart = shape.chart
                chart_info = {
                    "chartType": self._get_chart_type_name(chart.chart_type),
                    "hasLegend": chart.has_legend,
                }
                try:
                    chart_info["seriesCount"] = len(chart.series)
                except Exception:
                    pass
                slide_info["charts"].append(chart_info)
            except Exception as e:
                logger.warning(f"Error extracting chart: {e}")

        # Extract image info
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                slide_info["images"].append({
                    "width": shape.width,
                    "height": shape.height,
                })
        except Exception:
            pass

    def _extract_branding(self, prs: Presentation) -> Dict[str, Any]:
        """Extract branding information from a presentation."""
        branding = {
            "slideWidth": prs.slide_width,
            "slideHeight": prs.slide_height,
            "titleFontName": None,
            "bodyFontName": None,
            "titleFontSize": 28,
            "bodyFontSize": 12,
            "primaryColor": None,
            "secondaryColor": None,
            "accentColor": None,
        }

        try:
            slide_master = prs.slide_master
            if slide_master:
                try:
                    theme = slide_master.theme
                    if theme and theme.font_scheme:
                        branding["titleFontName"] = theme.font_scheme.major_font.typeface
                        branding["bodyFontName"] = theme.font_scheme.minor_font.typeface
                except Exception:
                    pass

                try:
                    theme = slide_master.theme
                    if theme and theme.color_scheme:
                        for color_name in ["accent1", "accent2", "dk1", "dk2"]:
                            try:
                                color = getattr(theme.color_scheme, color_name)
                                if color and hasattr(color, 'rgb') and color.rgb:
                                    hex_color = f"#{color.rgb}"
                                    if branding["primaryColor"] is None:
                                        branding["primaryColor"] = hex_color
                                    elif branding["secondaryColor"] is None:
                                        branding["secondaryColor"] = hex_color
                                    elif branding["accentColor"] is None:
                                        branding["accentColor"] = hex_color
                                        break
                            except Exception:
                                pass
                except Exception:
                    pass
        except Exception as e:
            logger.warning(f"Error extracting branding: {e}")

        branding["slideWidth"] = branding["slideWidth"] / 914400 if branding["slideWidth"] else 10
        branding["slideHeight"] = branding["slideHeight"] / 914400 if branding["slideHeight"] else 7.5

        return branding

    # =========================================================================
    # PPTX GENERATION
    # =========================================================================

    def _add_table_to_slide(
        self,
        slide,
        table_data: Dict[str, Any],
        top: int = None,
        slide_width: float = 10,
        font_name: Optional[str] = None,
        header_color=None,
    ) -> None:
        """Add a table to a slide with optional branding."""
        from pptx.enum.text import PP_ALIGN

        headers = table_data.get("headers", [])
        rows = table_data.get("rows", [])

        if not headers and not rows:
            return

        col_count = len(headers) if headers else (len(rows[0]) if rows else 0)
        row_count = (1 if headers else 0) + len(rows)

        if col_count == 0 or row_count == 0:
            return

        x = Inches(0.5)
        y = Inches(top) if top else Inches(1.5)
        width = Inches(slide_width - 1)
        height = Inches(0.4 * row_count)

        table_shape = slide.shapes.add_table(row_count, col_count, x, y, width, height)
        table = table_shape.table

        col_width = int(width / col_count)
        for col in table.columns:
            col.width = col_width

        row_idx = 0

        if headers:
            for col_idx, header in enumerate(headers):
                cell = table.cell(0, col_idx)
                cell.text = str(header) if header else ""
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.bold = True
                    paragraph.font.size = Pt(11)
                    paragraph.alignment = PP_ALIGN.CENTER
                    if font_name:
                        paragraph.font.name = font_name
                    if header_color:
                        paragraph.font.color.rgb = header_color
            row_idx = 1

        for data_row in rows:
            for col_idx, cell_value in enumerate(data_row):
                if col_idx < col_count:
                    cell = table.cell(row_idx, col_idx)
                    cell.text = str(cell_value) if cell_value else ""
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = Pt(10)
                        if font_name:
                            paragraph.font.name = font_name
            row_idx += 1

    def _add_bullet_list_to_slide(
        self,
        slide,
        items: List[str],
        top: int = None,
        font_name: Optional[str] = None,
        font_size: int = 12,
        bullet_color=None,
    ) -> None:
        """Add a bullet list to a slide with optional branding."""
        x = Inches(0.5)
        y = Inches(top) if top else Inches(1.5)
        width = Inches(9)
        height = Inches(0.3 * len(items))

        textbox = slide.shapes.add_textbox(x, y, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True

        for idx, item in enumerate(items):
            if idx == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(font_size)
            p.space_after = Pt(6)
            if font_name:
                p.font.name = font_name
            if bullet_color:
                p.font.color.rgb = bullet_color

    def _find_layout_by_name(self, prs: Presentation, layout_name: Optional[str]) -> Any:
        """Find a slide layout by name, with fallback to common layouts."""
        if not layout_name:
            layout_name = "Blank"

        layout_name_lower = layout_name.lower().strip()

        # Try exact match first
        for layout in prs.slide_layouts:
            if layout.name.lower().strip() == layout_name_lower:
                return layout

        # Try partial match
        for layout in prs.slide_layouts:
            if layout_name_lower in layout.name.lower():
                return layout

        # Common layout name mappings
        layout_mappings = {
            "title and content": ["title and content", "title_and_body", "title and body", "one_column_text", "content"],
            "title slide": ["title slide", "title"],
            "content": ["title and content", "title_and_body", "title and body", "one_column_text", "content", "two content"],
            "blank": ["blank"],
            "section": ["section header", "section", "section_header"],
        }

        for key, candidates in layout_mappings.items():
            if key in layout_name_lower:
                for candidate in candidates:
                    for layout in prs.slide_layouts:
                        if candidate in layout.name.lower():
                            return layout

        # Default to blank or first available layout
        try:
            return prs.slide_layouts[6]  # Usually blank
        except IndexError:
            return prs.slide_layouts[0]

    def _fill_slide_placeholders(
        self,
        slide,
        title: str = "",
        subtitle: str = "",
        body_text: List[str] = None,
    ) -> bool:
        """Fill placeholders in a slide. Returns True if any placeholder was filled."""
        from pptx.enum.shapes import PP_PLACEHOLDER

        filled = False
        body_text = body_text or []

        # Log available placeholders for debugging
        placeholder_types = []
        for shape in slide.shapes:
            if shape.is_placeholder:
                try:
                    ph_type = shape.placeholder_format.type
                    placeholder_types.append(f"{ph_type} ({shape.placeholder_format.idx})")
                except Exception as e:
                    logger.debug(f"Could not read placeholder type: {e}")
        logger.debug(f"Available placeholders: {placeholder_types}")

        for shape in slide.shapes:
            if not shape.is_placeholder:
                continue

            try:
                ph_type = shape.placeholder_format.type
            except Exception:
                continue

            # Handle title placeholders
            if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                if shape.has_text_frame:
                    tf = shape.text_frame
                    # Clear existing text first
                    for para in tf.paragraphs:
                        para.clear()
                    tf.paragraphs[0].text = title if title else ""
                    filled = True

            # Handle subtitle placeholders
            elif ph_type == PP_PLACEHOLDER.SUBTITLE:
                if shape.has_text_frame:
                    tf = shape.text_frame
                    for para in tf.paragraphs:
                        para.clear()
                    tf.paragraphs[0].text = subtitle if subtitle else ""
                    filled = True

            # Handle body/content placeholders
            # Note: Content placeholder in PowerPoint uses BODY type in python-pptx
            elif ph_type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
                if shape.has_text_frame:
                    tf = shape.text_frame
                    # Clear existing text first
                    for para in tf.paragraphs:
                        para.clear()

                    if body_text:
                        for i, text in enumerate(body_text):
                            if i == 0:
                                tf.paragraphs[0].text = text
                            else:
                                p = tf.add_paragraph()
                                p.text = text
                    else:
                        tf.paragraphs[0].text = ""
                    filled = True

            # Also handle generic placeholders that might contain text
            # Some templates use placeholder index rather than type
            elif shape.has_text_frame:
                # Check if this looks like a body placeholder by its index
                try:
                    idx = shape.placeholder_format.idx
                    # Index 1 is usually body/content in "Title and Content" layouts
                    if idx == 1 and body_text:
                        tf = shape.text_frame
                        for para in tf.paragraphs:
                            para.clear()
                        for i, text in enumerate(body_text):
                            if i == 0:
                                tf.paragraphs[0].text = text
                            else:
                                p = tf.add_paragraph()
                                p.text = text
                        filled = True
                except Exception:
                    pass

        return filled

    def _clone_slide(self, prs: Presentation, source_slide):
        """
        Clone a slide with all its shapes and relationships.

        This preserves decorative elements, images, and styling from the template.
        Based on: https://gist.github.com/robintw/3df1464e5c8a7ee8835e

        Important: We need to remove the default shapes that come from the layout
        before copying shapes from the source slide to avoid duplicates.

        Returns:
            pptx.slide.Slide: The cloned slide added to the presentation.
        """
        from copy import deepcopy
        from lxml import etree

        # Add a new slide with the same layout (this gives us the background)
        dest = prs.slides.add_slide(source_slide.slide_layout)

        # Get the shape tree (spTree) element
        spTree = dest.shapes._spTree

        # Remove all existing shapes from the new slide EXCEPT nvGrpSpPr and grpSpPr
        # These are the group properties that must remain
        shapes_to_remove = []
        for child in spTree:
            tag_name = etree.QName(child.tag).localname
            # Keep only the required group properties
            if tag_name not in ('nvGrpSpPr', 'grpSpPr'):
                shapes_to_remove.append(child)

        for shape_el in shapes_to_remove:
            spTree.remove(shape_el)

        # Now copy all shapes from source to destination
        for shp in source_slide.shapes:
            el = shp.element
            new_el = deepcopy(el)
            spTree.append(new_el)

        # Copy relationships (images, charts, etc.)
        for rel in source_slide.part.rels.values():
            # Skip notes slides
            if "notesSlide" in rel.reltype:
                continue
            # Add relationship to destination slide
            try:
                dest.part.rels.get_or_add(rel.reltype, rel._target)
            except Exception as e:
                # Log but continue - some relationships may not be copyable
                logger.debug(f"Could not copy relationship {rel.reltype}: {e}")

        return dest

    def _should_clone_slide(self, content: Dict[str, Any], idx: int) -> bool:
        """
        Determine if we should clone a template slide for this content.

        Only clone for title slides (title + subtitle, no body content).
        For content slides, use layout-based creation to avoid copying
        template's example content (like "01", "About us", etc.).
        """
        subtitle = content.get("subtitle", "")
        text_content = content.get("textContent", [])
        bullet_points = content.get("bulletPoints", [])
        tables = content.get("tables", [])
        has_body = bool(text_content or bullet_points or tables)

        # Only clone for title slides (subtitle present, no body)
        # or the first slide if it only has title
        if subtitle and not has_body:
            return True
        if idx == 0 and not has_body:
            return True

        return False

    def _find_template_slide_for_clone(
        self,
        template_slides: List[Any],
    ) -> Optional[Any]:
        """
        Find a template slide to clone (only for title slides).

        This is only called when _should_clone_slide returns True.
        """
        # For title slides, find a title slide layout
        preferred_layouts = ["title slide", "title", "section"]

        for preferred in preferred_layouts:
            for slide in template_slides:
                layout_name = slide.slide_layout.name.lower()
                if preferred in layout_name:
                    logger.debug(f"Found title slide to clone: {layout_name}")
                    return slide

        # Fallback to first slide if no title slide found
        if template_slides:
            return template_slides[0]

        return None

    def _extract_decorative_elements(self, source_slide) -> List[Any]:
        """
        Extract purely decorative elements from a template slide.

        Returns a list of XML elements that can be copied to a new slide.
        Only includes visual decorations (images, lines, shapes without text)
        and excludes any text-containing shapes to avoid copying template's
        example content.
        """
        from copy import deepcopy
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        decorative_elements = []

        for shape in source_slide.shapes:
            # Skip placeholders
            if shape.is_placeholder:
                continue

            # Skip any shape with text content (these are likely template examples)
            if shape.has_text_frame:
                text = self._get_shape_text(shape).strip()
                if text:
                    # Has text - skip this shape (likely "01", "About us", etc.)
                    continue

            # Include images
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    decorative_elements.append(deepcopy(shape.element))
                    continue
            except Exception:
                pass

            # Include lines and connectors
            try:
                if shape.shape_type in (MSO_SHAPE_TYPE.LINE, MSO_SHAPE_TYPE.FREEFORM):
                    decorative_elements.append(deepcopy(shape.element))
                    continue
            except Exception:
                pass

            # Include groups (often contain decorative patterns)
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    # Only include if no text in any child
                    has_text = False
                    for child_shape in shape.shapes:
                        if child_shape.has_text_frame and self._get_shape_text(child_shape).strip():
                            has_text = True
                            break
                    if not has_text:
                        decorative_elements.append(deepcopy(shape.element))
                    continue
            except Exception:
                pass

            # For other shapes without text, include them (borders, boxes, etc.)
            if not shape.has_text_frame:
                decorative_elements.append(deepcopy(shape.element))

        return decorative_elements

    def _get_best_layout_for_content(self, prs: Presentation, content: Dict[str, Any]) -> Any:
        """
        Get the best layout from the presentation's slide layouts based on content.

        This is used as a fallback when no matching template slide is found.
        """
        subtitle = content.get("subtitle", "")
        text_content = content.get("textContent", [])
        bullet_points = content.get("bulletPoints", [])
        tables = content.get("tables", [])
        has_body = bool(text_content or bullet_points or tables)

        if subtitle and not has_body:
            # Title + subtitle → title slide layout
            preferred = ["title slide", "title"]
        elif has_body:
            # Has body → content layout
            # Order matters: prefer layouts with body placeholder
            preferred = ["title and content", "title, content", "content", "title and body", "two content"]
        else:
            # Default → blank or content
            preferred = ["blank", "title and content"]

        # Log for debugging
        layout_names = [l.name.lower() for l in prs.slide_layouts]
        logger.debug(f"Looking for layout, preferred={preferred}, available={layout_names}")

        # Find matching layout
        for pref in preferred:
            for layout in prs.slide_layouts:
                if pref in layout.name.lower():
                    # Verify this layout has a body placeholder
                    if has_body:
                        has_body_placeholder = self._layout_has_body_placeholder(layout)
                        if has_body_placeholder:
                            return layout
                        else:
                            logger.debug(f"Layout '{layout.name}' has no body placeholder, skipping")
                            continue
                    return layout

        # Fallback: find any layout with a body placeholder
        if has_body:
            for layout in prs.slide_layouts:
                if self._layout_has_body_placeholder(layout):
                    logger.info(f"Using fallback layout with body placeholder: {layout.name}")
                    return layout

        # Last resort: first layout
        return prs.slide_layouts[0]

    def _layout_has_body_placeholder(self, layout) -> bool:
        """Check if a layout has a body/content placeholder."""
        from pptx.enum.shapes import PP_PLACEHOLDER

        for shape in layout.placeholders:
            try:
                ph_type = shape.placeholder_format.type
                if ph_type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
                    return True
            except Exception:
                pass
        return False

    def _create_slide_from_layout_only(
        self,
        prs: Presentation,
        content: Dict[str, Any],
        slide_width: float,
        title_font: Optional[str],
        body_font: Optional[str],
        title_font_size: int,
        body_font_size: int,
        primary_color,
        secondary_color,
    ) -> Any:
        """
        Create a new slide using only the template's layout.

        This method:
        1. Selects the best layout from the template's available layouts
        2. Creates a new slide with that layout (inherits theme styling)
        3. Fills placeholders with content

        The layout provides: background color/style, placeholder positions,
        theme fonts and colors. This avoids copying any template-specific
        content like "01", "About us", etc.
        """
        # Get the best layout for this content
        layout = self._get_best_layout_for_content(prs, content)

        # Log all available layouts for debugging
        all_layouts = [l.name for l in prs.slide_layouts]
        logger.info(f"Available layouts: {all_layouts}")
        logger.info(f"Selected layout: {layout.name}")

        # Create the new slide
        slide = prs.slides.add_slide(layout)

        # Log placeholder info for debugging
        for shape in slide.shapes:
            if shape.is_placeholder:
                try:
                    ph = shape.placeholder_format
                    logger.info(f"Placeholder: idx={ph.idx}, type={ph.type}")
                except Exception as e:
                    logger.warning(f"Error reading placeholder: {e}")

        # Fill placeholders with content
        title = content.get("title", "")
        subtitle = content.get("subtitle", "")
        text_content = content.get("textContent", [])
        bullet_points = content.get("bulletPoints", [])
        tables = content.get("tables", [])
        body_text = bullet_points if bullet_points else text_content

        # Try to fill placeholders first
        placeholders_filled = self._fill_slide_placeholders(slide, title, subtitle, body_text)

        # If placeholders couldn't be filled (e.g., blank layout), add content manually
        if not placeholders_filled:
            current_top = 0.5

            if title:
                self._add_title_to_slide(
                    slide, title, current_top, slide_width, title_font,
                    title_font_size, primary_color
                )
                current_top += 0.8

            if subtitle:
                self._add_subtitle_to_slide(
                    slide, subtitle, current_top, slide_width, body_font,
                    title_font_size - 10, secondary_color
                )
                current_top += 0.6

            for text in text_content:
                if isinstance(text, dict):
                    text = text.get("text", "")
                if text:
                    self._add_text_to_slide(slide, text, current_top, slide_width, body_font, body_font_size)
                    current_top += 0.3 + (len(text) // 100) * 0.2

            if bullet_points:
                self._add_bullet_list_to_slide(
                    slide, bullet_points, current_top, body_font, body_font_size, primary_color
                )
                current_top += 0.3 * len(bullet_points) + 0.3

        # Add tables (always added manually below placeholders)
        if tables:
            current_top = 4.0
            for table_data in tables:
                self._add_table_to_slide(
                    slide, table_data, current_top, slide_width, body_font, primary_color
                )
                row_count = len(table_data.get("rows", [])) + (1 if table_data.get("headers") else 0)
                current_top += 0.4 * row_count + 0.3

        return slide

    async def generate_pptx_from_content(
        self,
        slides_content: List[Dict[str, Any]],
        branding: Optional[Dict[str, Any]] = None,
        template_bytes: Optional[bytes] = None,
    ) -> bytes:
        """
        Generate a PPTX based on provided content.

        When template is provided, clones template slides to preserve decorative elements.
        Backend automatically selects the best template slide based on content fields:
        - title + subtitle (no body) → clone title slide
        - title + bulletPoints/textContent → clone content slide
        - title only → clone blank or first slide

        Args:
            slides_content: List of slide content dicts with title, subtitle, textContent, bulletPoints, tables
            branding: Optional branding info (used when no template provided)
            template_bytes: Optional template file bytes to use as base

        Returns:
            Generated PPTX as bytes
        """
        branding = branding or {}
        title_font = branding.get("titleFontName")
        body_font = branding.get("bodyFontName")
        title_font_size = branding.get("titleFontSize", 28)
        body_font_size = branding.get("bodyFontSize", 12)
        primary_color = self._parse_hex_color(branding.get("primaryColor"))
        secondary_color = self._parse_hex_color(branding.get("secondaryColor"))

        if template_bytes:
            logger.info("Using template file as base for PPTX generation")
            prs = Presentation(BytesIO(template_bytes))
            slide_width = prs.slide_width / 914400

            # Index template slides for cloning
            template_slides = list(prs.slides)
            original_slide_count = len(template_slides)
            logger.info(f"Template has {original_slide_count} slides")

            # Generate slides based on content type
            for idx, content in enumerate(slides_content):
                # Decide strategy: clone for title slides, layout-based for content slides
                if self._should_clone_slide(content, idx):
                    # Title slide - clone to preserve decorations
                    source_slide = self._find_template_slide_for_clone(template_slides)
                    if source_slide:
                        logger.info(f"Cloning title slide for content idx={idx}")
                        new_slide = self._clone_slide(prs, source_slide)
                        self._fill_cloned_slide(new_slide, content, slide_width, body_font, primary_color)
                    else:
                        # Fallback if no title slide found
                        logger.info(f"No title slide to clone for idx={idx}, using layout")
                        self._create_slide_from_layout_only(
                            prs, content, slide_width,
                            title_font, body_font, title_font_size, body_font_size,
                            primary_color, secondary_color
                        )
                else:
                    # Content slide - use layout only (no decoration copying)
                    # This avoids copying template's example content
                    logger.info(f"Creating content slide from layout for idx={idx}")
                    self._create_slide_from_layout_only(
                        prs, content, slide_width,
                        title_font, body_font, title_font_size, body_font_size,
                        primary_color, secondary_color
                    )

            # Remove original template slides (they are at the beginning)
            for _ in range(original_slide_count):
                rId = prs.slides._sldIdLst[0].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[0]

            logger.info(f"Removed {original_slide_count} original template slides")
        else:
            logger.info("Creating new blank PPTX (no template provided)")
            prs = Presentation()
            slide_width = branding.get("slideWidth", 10)
            slide_height = branding.get("slideHeight", 7.5)
            prs.slide_width = Inches(slide_width)
            prs.slide_height = Inches(slide_height)

            for idx, content in enumerate(slides_content):
                self._create_slide_from_layout(
                    prs, idx, content, slide_width,
                    title_font, body_font, title_font_size, body_font_size,
                    primary_color, secondary_color
                )

        output = BytesIO()
        prs.save(output)
        return output.getvalue()

    def _fill_cloned_slide(
        self,
        slide,
        content: Dict[str, Any],
        slide_width: float,
        body_font: Optional[str],
        primary_color,
    ) -> None:
        """
        Fill placeholders in a cloned slide with new content.

        When content is empty for a placeholder, the placeholder text is cleared
        to avoid showing "Click to add Text" messages.
        """
        from pptx.enum.shapes import PP_PLACEHOLDER

        title = content.get("title", "")
        subtitle = content.get("subtitle", "")
        text_content = content.get("textContent", [])
        bullet_points = content.get("bulletPoints", [])
        tables = content.get("tables", [])

        body_text = bullet_points if bullet_points else text_content

        # Track which placeholder types we've already processed to avoid duplicates
        # (cloned slides may have duplicate shapes from both layout and clone)
        processed_types = set()

        # Fill placeholders
        for shape in slide.shapes:
            if not shape.is_placeholder:
                continue

            try:
                ph_type = shape.placeholder_format.type
            except Exception:
                continue

            # Skip if we've already processed this placeholder type
            if ph_type in processed_types:
                continue

            if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                processed_types.add(ph_type)
                if shape.has_text_frame:
                    tf = shape.text_frame
                    # Clear all existing text
                    for para in tf.paragraphs:
                        para.clear()
                    # Set title if provided, otherwise leave empty (not "Click to add")
                    if title:
                        tf.paragraphs[0].text = title
                    else:
                        # Set empty string to clear default placeholder text
                        tf.paragraphs[0].text = ""

            elif ph_type == PP_PLACEHOLDER.SUBTITLE:
                processed_types.add(ph_type)
                if shape.has_text_frame:
                    tf = shape.text_frame
                    # Clear all existing text
                    for para in tf.paragraphs:
                        para.clear()
                    # Set subtitle if provided, otherwise leave empty
                    if subtitle:
                        tf.paragraphs[0].text = subtitle
                    else:
                        # Set empty string to clear default placeholder text
                        tf.paragraphs[0].text = ""

            elif ph_type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
                processed_types.add(ph_type)
                if shape.has_text_frame:
                    tf = shape.text_frame
                    # Clear existing paragraphs
                    for para in tf.paragraphs:
                        para.clear()

                    if body_text:
                        # Add new content
                        for i, text in enumerate(body_text):
                            if i == 0:
                                tf.paragraphs[0].text = text
                            else:
                                p = tf.add_paragraph()
                                p.text = text
                    else:
                        # Set empty string to clear default placeholder text
                        tf.paragraphs[0].text = ""

        # Add tables below placeholders
        if tables:
            current_top = 4.0  # Start below typical placeholder area
            for table_data in tables:
                self._add_table_to_slide(
                    slide, table_data, current_top, slide_width, body_font, primary_color
                )
                row_count = len(table_data.get("rows", [])) + (1 if table_data.get("headers") else 0)
                current_top += 0.4 * row_count + 0.3

    def _create_slide_from_layout(
        self,
        prs: Presentation,
        idx: int,
        content: Dict[str, Any],
        slide_width: float,
        title_font: Optional[str],
        body_font: Optional[str],
        title_font_size: int,
        body_font_size: int,
        primary_color,
        secondary_color,
    ):
        """Create a slide from layout (used when no template is provided)."""
        title = content.get("title", "")
        subtitle = content.get("subtitle", "")
        text_content = content.get("textContent", [])
        bullet_points = content.get("bulletPoints", [])
        tables = content.get("tables", [])

        # Use blank layout for new presentations
        try:
            layout = prs.slide_layouts[6]  # Usually blank
        except IndexError:
            layout = prs.slide_layouts[0]

        slide = prs.slides.add_slide(layout)

        # Add content manually
        current_top = 0.5

        if title:
            self._add_title_to_slide(
                slide, title, current_top, slide_width, title_font,
                title_font_size if idx == 0 else title_font_size - 4, primary_color
            )
            current_top += 0.8

        if subtitle:
            self._add_subtitle_to_slide(
                slide, subtitle, current_top, slide_width, body_font,
                title_font_size - 10 if idx == 0 else title_font_size - 12, secondary_color
            )
            current_top += 0.6

        for text in text_content:
            if isinstance(text, dict):
                text = text.get("text", "")
            if text:
                self._add_text_to_slide(slide, text, current_top, slide_width, body_font, body_font_size)
                current_top += 0.3 + (len(text) // 100) * 0.2

        if bullet_points:
            self._add_bullet_list_to_slide(
                slide, bullet_points, current_top, body_font, body_font_size, primary_color
            )
            current_top += 0.3 * len(bullet_points) + 0.3

        for table_data in tables:
            self._add_table_to_slide(
                slide, table_data, current_top, slide_width, body_font, primary_color
            )
            row_count = len(table_data.get("rows", [])) + (1 if table_data.get("headers") else 0)
            current_top += 0.4 * row_count + 0.3

        return slide

    def _add_title_to_slide(
        self, slide, title: str, top: float, slide_width: float,
        font_name: Optional[str], font_size: int, color
    ) -> None:
        """Add a title to a slide."""
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(top), Inches(slide_width - 1), Inches(0.6)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(font_size)
        p.font.bold = True
        if font_name:
            p.font.name = font_name
        if color:
            p.font.color.rgb = color

    def _add_subtitle_to_slide(
        self, slide, subtitle: str, top: float, slide_width: float,
        font_name: Optional[str], font_size: int, color
    ) -> None:
        """Add a subtitle to a slide."""
        sub_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(top), Inches(slide_width - 1), Inches(0.4)
        )
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(font_size)
        if font_name:
            p.font.name = font_name
        if color:
            p.font.color.rgb = color

    def _add_text_to_slide(
        self, slide, text: str, top: float, slide_width: float,
        font_name: Optional[str], font_size: int
    ) -> None:
        """Add text to a slide."""
        text_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(top), Inches(slide_width - 1), Inches(0.5)
        )
        tf = text_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = str(text)
        p.font.size = Pt(font_size)
        if font_name:
            p.font.name = font_name


# Global service instance
_pptx_service: Optional[PptxService] = None


def get_pptx_service() -> PptxService:
    """Get or create the global PPTX service instance."""
    global _pptx_service
    if _pptx_service is None:
        _pptx_service = PptxService()
    return _pptx_service
