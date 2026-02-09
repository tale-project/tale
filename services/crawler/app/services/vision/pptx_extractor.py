"""Smart PPTX text extraction with Vision API support.

This module extracts text and images from PPTX files while preserving
the relative positions of shapes within each slide.
"""

import asyncio
from io import BytesIO

from loguru import logger
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ...config import settings
from .openai_client import vision_client

MIN_IMAGE_SIZE = 10000  # ~100x100 pixels


async def _describe_image_bytes(
    image_bytes: bytes,
    semaphore: asyncio.Semaphore,
) -> str:
    """Describe image bytes using Vision API."""
    async with semaphore:
        try:
            if len(image_bytes) < MIN_IMAGE_SIZE:
                logger.debug(f"Skipping small image ({len(image_bytes)} bytes)")
                return ""
            return await vision_client.describe_image(image_bytes)
        except Exception as e:
            logger.warning(f"Failed to describe image: {e}")
            return ""


def _extract_table_text(table) -> str:
    """Extract text from a PPTX table."""
    rows_text = []
    for row in table.rows:
        cells_text = [cell.text.strip() for cell in row.cells]
        rows_text.append(" | ".join(cells_text))
    return "\n".join(rows_text)


async def _process_slide(
    slide_num: int,
    slide,
    semaphore: asyncio.Semaphore,
    process_images: bool,
) -> tuple[int, str, bool]:
    """Process a single slide and extract its content.

    Args:
        slide_num: Slide number (1-indexed)
        slide: PPTX slide object
        semaphore: Concurrency limiter
        process_images: Whether to process images

    Returns:
        Tuple of (slide_num, slide_content, vision_used)
    """
    elements: list[tuple[float, str]] = []  # (top, content)
    image_tasks: list[tuple[float, bytes]] = []  # (top, image_bytes)

    for shape in slide.shapes:
        top = shape.top or 0

        # Text frame
        if shape.has_text_frame:
            paragraphs_text = []
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    paragraphs_text.append(text)
            if paragraphs_text:
                elements.append((top, "\n".join(paragraphs_text)))

        # Table
        if shape.has_table:
            table_text = _extract_table_text(shape.table)
            if table_text:
                elements.append((top, f"[Table]\n{table_text}"))

        # Collect images for concurrent processing
        if process_images and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                image_tasks.append((top, shape.image.blob))
            except Exception as e:
                logger.warning(f"Failed to extract image on slide {slide_num}: {e}")

    # Process all images concurrently
    vision_used = False
    if image_tasks:
        results = await asyncio.gather(
            *[_describe_image_bytes(img_bytes, semaphore) for _, img_bytes in image_tasks],
            return_exceptions=True,
        )
        for (top, _), result in zip(image_tasks, results, strict=False):
            if isinstance(result, Exception):
                logger.warning(f"Failed to describe image on slide {slide_num}: {result}")
            elif result:
                elements.append((top, f"[Image: {result}]"))
                vision_used = True

    # Sort elements by top position
    elements.sort(key=lambda x: x[0])
    content = "\n\n".join(elem[1] for elem in elements)
    slide_text = f"--- Slide {slide_num} ---\n{content}"

    return slide_num, slide_text, vision_used


async def extract_text_from_pptx_bytes(
    pptx_bytes: bytes,
    filename: str = "presentation.pptx",
    *,
    process_images: bool = True,
) -> tuple[list[str], bool]:
    """Extract text from PPTX bytes with Vision support.

    Args:
        pptx_bytes: Raw PPTX bytes
        filename: Optional filename for logging
        process_images: Whether to extract and describe embedded images

    Returns:
        Tuple of (list of slide contents, vision_used flag)
    """
    logger.info(f"Processing PPTX: {filename}")

    prs = Presentation(BytesIO(pptx_bytes))
    semaphore = asyncio.Semaphore(settings.vision_max_concurrent_pages)

    # Process all slides concurrently
    tasks = [_process_slide(i + 1, slide, semaphore, process_images) for i, slide in enumerate(prs.slides)]
    results = await asyncio.gather(*tasks, return_exceptions=True)

    slides_content: list[tuple[int, str]] = []
    vision_used = False

    for result in results:
        if isinstance(result, Exception):
            logger.warning(f"Slide processing failed: {result}")
            continue
        slide_num, content, slide_vision_used = result
        slides_content.append((slide_num, content))
        if slide_vision_used:
            vision_used = True

    # Sort by slide number
    slides_content.sort(key=lambda x: x[0])
    ordered_content = [s[1] for s in slides_content]

    logger.info(f"PPTX processing complete: {len(ordered_content)} slides, Vision API used: {vision_used}")

    return ordered_content, vision_used
