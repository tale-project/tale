"""
File Parser Service for extracting text content from documents.

Handles:
- PDF text extraction using PyMuPDF (with optional Vision API for images/OCR)
- DOCX text extraction using python-docx (with optional Vision API for images)
- PPTX text extraction using python-pptx (with optional Vision API for images)
"""

import datetime as dt
import logging
import re
from io import BytesIO
from typing import Any

from ..config import settings

logger = logging.getLogger(__name__)

_PDF_DATE_RE = re.compile(
    r"^(?:D:)?"
    r"(\d{4})"
    r"(\d{2})?"
    r"(\d{2})?"
    r"(\d{2})?"
    r"(\d{2})?"
    r"(\d{2})?"
    r"([Z+\-])?"
    r"(\d{2})?'?"
    r"(\d{2})?'?"
)

_MIN_YEAR = 1970
_MAX_YEAR = 2100


def _parse_pdf_date(date_str: str | None) -> int | None:
    """Parse PDF date format ``D:YYYYMMDDHHmmSSOHH'mm'`` to Unix timestamp ms.

    Returns ``None`` for missing, malformed, or out-of-range dates.
    """
    if not date_str or not isinstance(date_str, str):
        return None

    date_str = date_str.strip()
    if not date_str:
        return None

    match = _PDF_DATE_RE.match(date_str)
    if not match:
        logger.warning("Failed to parse PDF date: '%s'", date_str)
        return None

    try:
        year = int(match.group(1))
        if year < _MIN_YEAR or year > _MAX_YEAR:
            logger.warning("PDF date year out of range (%d): '%s'", year, date_str)
            return None

        month = int(match.group(2) or "01")
        day = int(match.group(3) or "01")
        hour = int(match.group(4) or "00")
        minute = int(match.group(5) or "00")
        second = int(match.group(6) or "00")

        tz_sign = match.group(7)
        tz_hours = int(match.group(8) or "0")
        tz_minutes = int(match.group(9) or "0")

        if tz_sign == "-":
            tz_offset = dt.timezone(dt.timedelta(hours=-tz_hours, minutes=-tz_minutes))
        elif tz_sign == "+":
            tz_offset = dt.timezone(dt.timedelta(hours=tz_hours, minutes=tz_minutes))
        else:
            tz_offset = dt.UTC

        d = dt.datetime(year, month, day, hour, minute, second, tzinfo=tz_offset)
        return int(d.timestamp() * 1000)
    except (ValueError, OverflowError) as e:
        logger.warning("Failed to parse PDF date '%s': %s", date_str, e)
        return None


def _to_unix_ms(d: dt.datetime | None) -> int | None:
    """Convert a datetime to Unix timestamp in milliseconds.

    Handles both timezone-aware and naive (assumed UTC) datetimes.
    Returns ``None`` for ``None`` input.
    """
    if d is None:
        return None
    if not isinstance(d, dt.datetime):
        return None
    try:
        if d.tzinfo is None:
            d = d.replace(tzinfo=dt.UTC)
        year = d.year
        if year < _MIN_YEAR or year > _MAX_YEAR:
            logger.warning("Datetime year out of range (%d)", year)
            return None
        return int(d.timestamp() * 1000)
    except (ValueError, OverflowError, OSError) as e:
        logger.warning("Failed to convert datetime to Unix ms: %s", e)
        return None


def _extract_ooxml_metadata(file_bytes: bytes, fmt: str) -> dict[str, Any]:
    """Extract dates + title/author from OOXML (DOCX/PPTX) core properties."""
    result: dict[str, Any] = {
        "created_at": None,
        "modified_at": None,
        "title": "",
        "author": "",
    }

    try:
        if fmt == "docx":
            from docx import Document

            doc = Document(BytesIO(file_bytes))
            props = doc.core_properties
        elif fmt == "pptx":
            from pptx import Presentation

            prs = Presentation(BytesIO(file_bytes))
            props = prs.core_properties
        else:
            return result

        result["title"] = props.title or ""
        result["author"] = getattr(props, "author", "") or ""
        result["created_at"] = _to_unix_ms(props.created)
        result["modified_at"] = _to_unix_ms(props.modified)

    except Exception as e:
        logger.warning("Failed to extract OOXML metadata from %s: %s", fmt, e)

    return result


class FileParserService:
    """Service for parsing and extracting text from various document formats."""

    def parse_pdf(self, file_bytes: bytes, filename: str = "document.pdf") -> dict[str, Any]:
        """Extract text content from a PDF file (basic extraction without Vision)."""
        import fitz  # PyMuPDF

        try:
            doc = fitz.open(stream=file_bytes, filetype="pdf")
            pages = []
            full_text = []

            for page_num, page in enumerate(doc, start=1):
                text = page.get_text("text")
                pages.append({"page_number": page_num, "text": text.strip()})
                full_text.append(text)

            metadata = doc.metadata or {}
            doc.close()

            return {
                "success": True,
                "filename": filename,
                "file_type": "application/pdf",
                "page_count": len(pages),
                "pages": pages,
                "full_text": "\n\n".join(full_text).strip(),
                "metadata": {
                    "title": metadata.get("title", ""),
                    "author": metadata.get("author", ""),
                    "subject": metadata.get("subject", ""),
                    "created_at": _parse_pdf_date(metadata.get("creationDate")),
                    "modified_at": _parse_pdf_date(metadata.get("modDate")),
                },
            }
        except Exception as e:
            logger.error(f"Error parsing PDF: {e}")
            return {"success": False, "filename": filename, "file_type": "application/pdf", "error": str(e)}

    async def parse_pdf_with_vision(
        self,
        file_bytes: bytes,
        filename: str = "document.pdf",
        user_input: str | None = None,
        process_images: bool = True,
        ocr_scanned_pages: bool = True,
        model: str | None = None,
    ) -> dict[str, Any]:
        """Extract text content from a PDF file using Vision API.

        Args:
            file_bytes: Raw PDF bytes
            filename: Filename for logging
            user_input: Optional user instruction for AI extraction
            process_images: Whether to extract and describe embedded images
            ocr_scanned_pages: Whether to OCR pages with low text content

        Returns:
            Extraction result with full_text and metadata
        """
        from .vision.openai_client import UsageAccumulator, process_pages_with_llm
        from .vision.pdf_extractor import extract_text_from_pdf_bytes

        try:
            acc = UsageAccumulator()

            pages_content, vision_used = await extract_text_from_pdf_bytes(
                file_bytes,
                filename,
                process_images=process_images,
                ocr_scanned_pages=ocr_scanned_pages,
                usage=acc,
            )

            resolved_model = model
            if user_input:
                pages_content = await process_pages_with_llm(
                    pages_content,
                    user_input,
                    max_concurrent=3,
                    model=model,
                    usage=acc,
                )
                resolved_model = model or settings.get_fast_model()

            import fitz as _fitz

            pdf_metadata: dict[str, Any] = {}
            try:
                with _fitz.open(stream=file_bytes, filetype="pdf") as _doc:
                    _raw = _doc.metadata or {}
                    pdf_metadata = {
                        "title": _raw.get("title", ""),
                        "author": _raw.get("author", ""),
                        "created_at": _parse_pdf_date(_raw.get("creationDate")),
                        "modified_at": _parse_pdf_date(_raw.get("modDate")),
                    }
            except Exception as e:
                logger.warning("Failed to extract PDF vision metadata: %s", e)

            result: dict[str, Any] = {
                "success": True,
                "filename": filename,
                "file_type": "application/pdf",
                "page_count": len(pages_content),
                "full_text": "\n\n".join(pages_content),
                "vision_used": vision_used,
                "metadata": pdf_metadata,
            }
            if acc.total_tokens > 0:
                result["usage"] = acc.to_dict(model=resolved_model)
            return result
        except Exception as e:
            logger.error(f"Error parsing PDF with Vision: {e}")
            return {
                "success": False,
                "filename": filename,
                "file_type": "application/pdf",
                "error": str(e),
            }

    def parse_docx(self, file_bytes: bytes, filename: str = "document.docx") -> dict[str, Any]:
        """Extract text content from a DOCX file (basic extraction without Vision)."""
        from docx import Document

        try:
            doc = Document(BytesIO(file_bytes))
            paragraphs = []
            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    paragraphs.append({"text": text, "style": para.style.name if para.style else None})

            tables = []
            for table in doc.tables:
                table_data = [[cell.text.strip() for cell in row.cells] for row in table.rows]
                if table_data:
                    tables.append(table_data)

            full_text = "\n".join(p["text"] for p in paragraphs)
            core_props = doc.core_properties

            return {
                "success": True,
                "filename": filename,
                "file_type": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                "paragraph_count": len(paragraphs),
                "table_count": len(tables),
                "paragraphs": paragraphs,
                "tables": tables,
                "full_text": full_text,
                "metadata": {
                    "title": core_props.title or "",
                    "author": core_props.author or "",
                    "created_at": _to_unix_ms(core_props.created) if hasattr(core_props, "created") else None,
                    "modified_at": _to_unix_ms(core_props.modified) if hasattr(core_props, "modified") else None,
                },
            }
        except Exception as e:
            logger.error(f"Error parsing DOCX: {e}")
            return {"success": False, "filename": filename, "error": str(e)}

    async def parse_docx_with_vision(
        self,
        file_bytes: bytes,
        filename: str = "document.docx",
        user_input: str | None = None,
        process_images: bool = True,
        model: str | None = None,
    ) -> dict[str, Any]:
        """Extract text content from a DOCX file using Vision API.

        Args:
            file_bytes: Raw DOCX bytes
            filename: Filename for logging
            user_input: Optional user instruction for AI extraction
            process_images: Whether to extract and describe embedded images

        Returns:
            Extraction result with full_text and metadata
        """
        from .vision.docx_extractor import extract_text_from_docx_bytes
        from .vision.openai_client import UsageAccumulator, process_pages_with_llm

        try:
            acc = UsageAccumulator()

            content_list, vision_used = await extract_text_from_docx_bytes(
                file_bytes,
                filename,
                process_images=process_images,
                usage=acc,
            )

            resolved_model = model
            if user_input:
                content_list = await process_pages_with_llm(
                    content_list,
                    user_input,
                    max_concurrent=3,
                    model=model,
                    usage=acc,
                )
                resolved_model = model or settings.get_fast_model()

            docx_dates = _extract_ooxml_metadata(file_bytes, "docx")

            result: dict[str, Any] = {
                "success": True,
                "filename": filename,
                "file_type": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                "element_count": len(content_list),
                "full_text": "\n\n".join(content_list),
                "vision_used": vision_used,
                "metadata": docx_dates,
            }
            if acc.total_tokens > 0:
                result["usage"] = acc.to_dict(model=resolved_model)
            return result
        except Exception as e:
            logger.error(f"Error parsing DOCX with Vision: {e}")
            return {
                "success": False,
                "filename": filename,
                "error": str(e),
            }

    def parse_pptx(self, file_bytes: bytes, filename: str = "presentation.pptx") -> dict[str, Any]:
        """Extract text content from a PPTX file (basic extraction without Vision)."""
        from pptx import Presentation

        try:
            prs = Presentation(BytesIO(file_bytes))
            slides = []
            full_text_parts = []

            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_text = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                slide_text.append(text)
                    if shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                text = cell.text.strip()
                                if text:
                                    slide_text.append(text)

                slides.append(
                    {
                        "slide_number": slide_num,
                        "text_content": slide_text,
                        "full_text": "\n".join(slide_text),
                    }
                )
                full_text_parts.extend(slide_text)

            core_props = prs.core_properties
            return {
                "success": True,
                "filename": filename,
                "file_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                "slide_count": len(slides),
                "slides": slides,
                "full_text": "\n\n".join(full_text_parts),
                "metadata": {
                    "title": core_props.title or "",
                    "author": core_props.author or "",
                    "created_at": _to_unix_ms(core_props.created) if hasattr(core_props, "created") else None,
                    "modified_at": _to_unix_ms(core_props.modified) if hasattr(core_props, "modified") else None,
                },
            }
        except Exception as e:
            logger.error(f"Error parsing PPTX: {e}")
            return {"success": False, "filename": filename, "error": str(e)}

    async def parse_pptx_with_vision(
        self,
        file_bytes: bytes,
        filename: str = "presentation.pptx",
        user_input: str | None = None,
        process_images: bool = True,
        model: str | None = None,
    ) -> dict[str, Any]:
        """Extract text content from a PPTX file using Vision API.

        Args:
            file_bytes: Raw PPTX bytes
            filename: Filename for logging
            user_input: Optional user instruction for AI extraction
            process_images: Whether to extract and describe embedded images

        Returns:
            Extraction result with full_text and metadata
        """
        from .vision.openai_client import UsageAccumulator, process_pages_with_llm
        from .vision.pptx_extractor import extract_text_from_pptx_bytes

        try:
            acc = UsageAccumulator()

            slides_content, vision_used = await extract_text_from_pptx_bytes(
                file_bytes,
                filename,
                process_images=process_images,
                usage=acc,
            )

            resolved_model = model
            if user_input:
                slides_content = await process_pages_with_llm(
                    slides_content,
                    user_input,
                    max_concurrent=3,
                    model=model,
                    usage=acc,
                )
                resolved_model = model or settings.get_fast_model()

            pptx_dates = _extract_ooxml_metadata(file_bytes, "pptx")

            result: dict[str, Any] = {
                "success": True,
                "filename": filename,
                "file_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                "slide_count": len(slides_content),
                "full_text": "\n\n".join(slides_content),
                "vision_used": vision_used,
                "metadata": pptx_dates,
            }
            if acc.total_tokens > 0:
                result["usage"] = acc.to_dict(model=resolved_model)
            return result
        except Exception as e:
            logger.error(f"Error parsing PPTX with Vision: {e}")
            return {
                "success": False,
                "filename": filename,
                "error": str(e),
            }

    def parse_file(self, file_bytes: bytes, filename: str, content_type: str = "") -> dict[str, Any]:
        """Parse a file based on its content type or filename extension (basic, without Vision)."""
        filename_lower = filename.lower()
        content_type_lower = content_type.lower() if content_type else ""

        if filename_lower.endswith(".pdf") or "pdf" in content_type_lower:
            return self.parse_pdf(file_bytes, filename)
        elif filename_lower.endswith(".docx") or "wordprocessingml" in content_type_lower:
            return self.parse_docx(file_bytes, filename)
        elif filename_lower.endswith(".pptx") or "presentationml" in content_type_lower:
            return self.parse_pptx(file_bytes, filename)
        else:
            return {
                "success": False,
                "filename": filename,
                "error": f"Unsupported file type: {filename} ({content_type}). Supported: PDF, DOCX, PPTX.",
            }

    async def parse_file_with_vision(
        self,
        file_bytes: bytes,
        filename: str,
        content_type: str = "",
        user_input: str | None = None,
        process_images: bool = True,
        ocr_scanned_pages: bool = True,
        model: str | None = None,
    ) -> dict[str, Any]:
        """Parse a file with Vision API support.

        Args:
            file_bytes: Raw file bytes
            filename: Filename
            content_type: MIME content type
            user_input: Optional user instruction for AI extraction
            process_images: Whether to extract and describe embedded images
            ocr_scanned_pages: Whether to OCR pages with low text content (PDF only)

        Returns:
            Extraction result with full_text and metadata
        """
        filename_lower = filename.lower()
        content_type_lower = content_type.lower() if content_type else ""

        if filename_lower.endswith(".pdf") or "pdf" in content_type_lower:
            return await self.parse_pdf_with_vision(
                file_bytes, filename, user_input, process_images, ocr_scanned_pages, model=model
            )
        elif filename_lower.endswith(".docx") or "wordprocessingml" in content_type_lower:
            return await self.parse_docx_with_vision(file_bytes, filename, user_input, process_images, model=model)
        elif filename_lower.endswith(".pptx") or "presentationml" in content_type_lower:
            return await self.parse_pptx_with_vision(file_bytes, filename, user_input, process_images, model=model)
        else:
            return {
                "success": False,
                "filename": filename,
                "error": f"Unsupported file type: {filename} ({content_type}). Supported: PDF, DOCX, PPTX.",
            }


_file_parser_service: FileParserService | None = None


def get_file_parser_service() -> FileParserService:
    """Get or create the file parser service instance."""
    global _file_parser_service
    if _file_parser_service is None:
        _file_parser_service = FileParserService()
    return _file_parser_service
