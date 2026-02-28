"""Tests for PPTX extraction."""

from __future__ import annotations

from io import BytesIO
from unittest.mock import MagicMock, PropertyMock

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from tale_knowledge.extraction.pptx import (
    _iter_shapes,
    _process_slide,
    extract_text_from_pptx_bytes,
)


def _make_simple_pptx(text: str = "Hello World") -> bytes:
    """Create a minimal PPTX with a single text slide."""
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]  # blank layout
    slide = prs.slides.add_slide(slide_layout)
    txbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    txbox.text_frame.text = text
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_grouped_pptx() -> bytes:
    """Create a PPTX with shapes inside a group.

    python-pptx doesn't expose a public API for creating group shapes,
    so we mock shapes directly in the tests instead. This helper creates
    a normal PPTX that we use for integration-level tests only.
    """
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    txbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    txbox.text_frame.text = "Top-level shape"
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


class TestIterShapes:
    def test_flat_shapes_returned_as_is(self):
        shape1 = MagicMock()
        shape1.shape_type = MSO_SHAPE_TYPE.TEXT_BOX
        shape2 = MagicMock()
        shape2.shape_type = MSO_SHAPE_TYPE.PICTURE

        result = list(_iter_shapes([shape1, shape2]))
        assert result == [shape1, shape2]

    def test_group_shapes_recursed(self):
        inner_shape = MagicMock()
        inner_shape.shape_type = MSO_SHAPE_TYPE.TEXT_BOX

        group = MagicMock()
        group.shape_type = MSO_SHAPE_TYPE.GROUP
        group.shapes = [inner_shape]

        result = list(_iter_shapes([group]))
        assert result == [inner_shape]

    def test_nested_groups(self):
        deepest = MagicMock()
        deepest.shape_type = MSO_SHAPE_TYPE.TEXT_BOX

        inner_group = MagicMock()
        inner_group.shape_type = MSO_SHAPE_TYPE.GROUP
        inner_group.shapes = [deepest]

        outer_group = MagicMock()
        outer_group.shape_type = MSO_SHAPE_TYPE.GROUP
        outer_group.shapes = [inner_group]

        result = list(_iter_shapes([outer_group]))
        assert result == [deepest]

    def test_mixed_flat_and_grouped(self):
        flat = MagicMock()
        flat.shape_type = MSO_SHAPE_TYPE.TEXT_BOX

        inner = MagicMock()
        inner.shape_type = MSO_SHAPE_TYPE.PICTURE

        group = MagicMock()
        group.shape_type = MSO_SHAPE_TYPE.GROUP
        group.shapes = [inner]

        result = list(_iter_shapes([flat, group]))
        assert result == [flat, inner]

    def test_empty_group(self):
        group = MagicMock()
        group.shape_type = MSO_SHAPE_TYPE.GROUP
        group.shapes = []

        result = list(_iter_shapes([group]))
        assert result == []


class TestExtractTextFromPptxBytes:
    @pytest.mark.asyncio
    async def test_basic_text_extraction(self):
        pptx_bytes = _make_simple_pptx("Hello World")
        text, vision_used = await extract_text_from_pptx_bytes(pptx_bytes)
        assert "Hello World" in text
        assert vision_used is False

    @pytest.mark.asyncio
    async def test_multiple_slides(self):
        prs = Presentation()
        for i in range(3):
            slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(slide_layout)
            txbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
            txbox.text_frame.text = f"Slide {i + 1} content"
        buf = BytesIO()
        prs.save(buf)

        text, _ = await extract_text_from_pptx_bytes(buf.getvalue())
        assert "Slide 1 content" in text
        assert "Slide 2 content" in text
        assert "Slide 3 content" in text

    @pytest.mark.asyncio
    async def test_slide_ordering(self):
        prs = Presentation()
        for i in range(3):
            slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(slide_layout)
            txbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
            txbox.text_frame.text = f"Slide {i + 1}"
        buf = BytesIO()
        prs.save(buf)

        text, _ = await extract_text_from_pptx_bytes(buf.getvalue())
        pos1 = text.index("Slide 1")
        pos2 = text.index("Slide 2")
        pos3 = text.index("Slide 3")
        assert pos1 < pos2 < pos3

    @pytest.mark.asyncio
    async def test_no_vision_without_client(self):
        pptx_bytes = _make_simple_pptx("No vision")
        text, vision_used = await extract_text_from_pptx_bytes(
            pptx_bytes, vision_client=None
        )
        assert "No vision" in text
        assert vision_used is False

    @pytest.mark.asyncio
    async def test_invalid_file_raises(self):
        with pytest.raises(ValueError, match="Invalid or corrupt file"):
            await extract_text_from_pptx_bytes(b"not a pptx file")

    @pytest.mark.asyncio
    async def test_table_extraction(self):
        prs = Presentation()
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        table_shape = slide.shapes.add_table(
            2, 2, Inches(1), Inches(1), Inches(4), Inches(2)
        )
        table = table_shape.table
        table.cell(0, 0).text = "A1"
        table.cell(0, 1).text = "B1"
        table.cell(1, 0).text = "A2"
        table.cell(1, 1).text = "B2"
        buf = BytesIO()
        prs.save(buf)

        text, _ = await extract_text_from_pptx_bytes(buf.getvalue())
        assert "A1" in text
        assert "[Table]" in text


class TestProcessSlideWithGroupedShapes:
    """Test that _process_slide extracts content from grouped shapes."""

    @pytest.mark.asyncio
    async def test_text_inside_group_extracted(self):
        import asyncio

        text_frame = MagicMock()
        text_frame.paragraphs = [MagicMock(text="Grouped text")]

        inner_shape = MagicMock()
        inner_shape.top = 100
        inner_shape.shape_type = MSO_SHAPE_TYPE.TEXT_BOX
        inner_shape.has_text_frame = True
        inner_shape.text_frame = text_frame
        inner_shape.has_table = False

        group = MagicMock()
        group.shape_type = MSO_SHAPE_TYPE.GROUP
        group.shapes = [inner_shape]

        slide = MagicMock()
        slide.shapes = [group]
        slide.has_notes_slide = False

        sem = asyncio.Semaphore(3)
        _, content, _ = await _process_slide(1, slide, sem, None, False)
        assert "Grouped text" in content

    @pytest.mark.asyncio
    async def test_table_inside_group_extracted(self):
        import asyncio

        row_mock = MagicMock()
        row_mock.cells = [MagicMock(text="Cell1"), MagicMock(text="Cell2")]

        table_mock = MagicMock()
        table_mock.rows = [row_mock]

        inner_shape = MagicMock()
        inner_shape.top = 200
        inner_shape.shape_type = MSO_SHAPE_TYPE.AUTO_SHAPE
        inner_shape.has_text_frame = False
        inner_shape.has_table = True
        inner_shape.table = table_mock

        group = MagicMock()
        group.shape_type = MSO_SHAPE_TYPE.GROUP
        group.shapes = [inner_shape]

        slide = MagicMock()
        slide.shapes = [group]
        slide.has_notes_slide = False

        sem = asyncio.Semaphore(3)
        _, content, _ = await _process_slide(1, slide, sem, None, False)
        assert "[Table]" in content
        assert "Cell1" in content
