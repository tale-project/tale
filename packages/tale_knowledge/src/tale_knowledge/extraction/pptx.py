"""Smart PPTX text extraction with Vision API support.

Extracts text and images from PPTX files while preserving
the relative positions of shapes within each slide.
"""

from __future__ import annotations

import asyncio
import zipfile
from io import BytesIO
from typing import TYPE_CHECKING

from loguru import logger
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ._helpers import describe_image_bytes, extract_table_text

if TYPE_CHECKING:
    from tale_knowledge.vision.client import VisionClient

MAX_UNCOMPRESSED_SIZE = 500 * 1024 * 1024  # 500 MB


async def _process_slide(
    slide_num: int,
    slide,
    semaphore: asyncio.Semaphore,
    vision_client: VisionClient | None,
    process_images: bool,
) -> tuple[int, str, bool]:
    elements: list[tuple[float, str]] = []
    image_tasks: list[tuple[float, bytes]] = []

    for shape in slide.shapes:
        top = shape.top or 0

        if shape.has_text_frame:
            paragraphs_text = []
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    paragraphs_text.append(text)
            if paragraphs_text:
                elements.append((top, "\n".join(paragraphs_text)))

        if shape.has_table:
            table_text = extract_table_text(shape.table)
            if table_text:
                elements.append((top, f"[Table]\n{table_text}"))

        if (
            process_images
            and vision_client
            and shape.shape_type == MSO_SHAPE_TYPE.PICTURE
        ):
            try:
                image_tasks.append((top, shape.image.blob))
            except Exception as e:
                logger.warning(f"Failed to extract image on slide {slide_num}: {e}")

    vision_used = False
    if image_tasks and vision_client:
        results = await asyncio.gather(
            *[
                describe_image_bytes(img_bytes, semaphore, vision_client)
                for _, img_bytes in image_tasks
            ],
            return_exceptions=True,
        )
        for (top, _), result in zip(image_tasks, results, strict=False):
            if isinstance(result, Exception):
                logger.warning(
                    f"Failed to describe image on slide {slide_num}: {result}"
                )
            elif result:
                elements.append((top, f"[Image: {result}]"))
                vision_used = True

    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
        notes_text = slide.notes_slide.notes_text_frame.text.strip()
        if notes_text:
            elements.append((float("inf"), f"[Notes]\n{notes_text}"))

    elements.sort(key=lambda x: x[0])
    content = "\n\n".join(elem[1] for elem in elements)
    slide_text = f"--- Slide {slide_num} ---\n{content}"

    return slide_num, slide_text, vision_used


async def extract_text_from_pptx_bytes(
    pptx_bytes: bytes,
    filename: str = "presentation.pptx",
    *,
    vision_client: VisionClient | None = None,
    process_images: bool = True,
    max_concurrent: int = 3,
) -> tuple[str, bool]:
    """Extract text from PPTX bytes with optional Vision support.

    Args:
        pptx_bytes: Raw PPTX bytes.
        filename: Filename for logging.
        vision_client: Optional VisionClient for image description.
        process_images: Whether to extract and describe embedded images.
        max_concurrent: Max concurrent Vision API calls.

    Returns:
        Tuple of (extracted_text, vision_was_used).
    """
    logger.info(f"Processing PPTX: {filename}")

    try:
        zf = zipfile.ZipFile(BytesIO(pptx_bytes))
        total = sum(info.file_size for info in zf.infolist())
        if total > MAX_UNCOMPRESSED_SIZE:
            raise ValueError(f"File exceeds maximum decompressed size ({total} bytes)")
        zf.close()
    except zipfile.BadZipFile:
        raise ValueError("Invalid or corrupt file")

    prs = Presentation(BytesIO(pptx_bytes))
    semaphore = asyncio.Semaphore(max_concurrent)

    tasks = [
        _process_slide(i + 1, slide, semaphore, vision_client, process_images)
        for i, slide in enumerate(prs.slides)
    ]
    results = await asyncio.gather(*tasks, return_exceptions=True)

    slides_content: list[tuple[int, str]] = []
    vision_used = False

    for result in results:
        if isinstance(result, Exception):
            logger.warning(f"Slide processing failed: {result}")
            continue
        slide_num, content, slide_vision_used = result
        slides_content.append((slide_num, content))
        if slide_vision_used:
            vision_used = True

    slides_content.sort(key=lambda x: x[0])
    combined_text = "\n\n".join(s[1] for s in slides_content)

    logger.info(
        f"PPTX processing complete: {len(slides_content)} slides, Vision API used: {vision_used}"
    )

    return combined_text, vision_used
